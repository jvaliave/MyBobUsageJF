#!/usr/bin/env python3
"""
Create PowerPoint presentation from IKC and Lineage Epic Analysis Reports
Enhanced version with complete data display
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime
import os
import math

def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section_header(prs, title):
    """Add section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    return slide

def add_content_slide(prs, title):
    """Add content slide with title."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    return slide

def add_chart_to_slide(slide, chart_data, chart_type, left, top, width, height, title=None):
    """Add chart to slide."""
    chart = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    ).chart
    
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    return chart

def add_table_slide(prs, title, data, col_widths=None):
    """Add a slide with a table."""
    slide = add_content_slide(prs, title)
    
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if cols == 0:
        return slide
    
    # Calculate table dimensions
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths if provided
    if col_widths:
        for i, col_width in enumerate(col_widths):
            table.columns[i].width = Inches(col_width)
    
    # Fill table with data
    for i, row_data in enumerate(data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_data)
            
            # Style header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.size = Pt(11)
            else:
                # Alternate row colors
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(9)
    
    return slide

def add_blocked_epics_table(prs, title, blocked_df):
    """Add blocked epics table with red styling."""
    slide = add_content_slide(prs, title)
    
    if blocked_df.empty:
        return slide
    
    # Prepare table data
    table_data = [['Epic Ref', 'Epic Name', 'Epic URL', 'Git URL', 'Epic Tags']]
    
    for idx, row in blocked_df.iterrows():
        epic_ref = str(row.get('Epic reference #', 'N/A'))
        epic_name = str(row.get('Epic name', 'N/A'))
        epic_url = str(row.get('Epic URL', 'N/A'))
        git_url = str(row.get('Github enterprise html_url', 'N/A'))
        epic_tags = str(row.get('Epic tags', 'N/A'))
        
        table_data.append([epic_ref, epic_name, epic_url, git_url, epic_tags])
    
    rows = len(table_data)
    cols = 5
    
    # Create table
    left = Inches(0.3)
    top = Inches(1.8)
    width = Inches(9.4)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(1.2)  # Epic Ref
    table.columns[1].width = Inches(2.0)  # Epic Name
    table.columns[2].width = Inches(2.5)  # Epic URL
    table.columns[3].width = Inches(2.5)  # Git URL
    table.columns[4].width = Inches(1.2)  # Epic Tags
    
    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_data)
            cell.text_frame.word_wrap = True
            
            # Style header row with red
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 107, 107)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.size = Pt(10)
            else:
                # Alternate row colors with light red
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 229, 229)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)
    
    return slide

def analyze_data(filename):
    """Analyze epic data from Excel file."""
    df = pd.read_excel(filename)
    
    # Status analysis - keep full status text
    status_counts = df['Epic status'].value_counts()
    
    # Tags analysis - get ALL tags
    all_tags = []
    for tags in df['Epic tags'].dropna():
        tag_list = [t.strip() for t in str(tags).split(',')]
        all_tags.extend(tag_list)
    
    tag_counts = pd.Series(all_tags).value_counts()
    
    # Company association analysis
    company_df = df[df['Company Association'].notna()].copy()
    company_count = len(company_df)
    
    # Blocked epics
    blocked_df = df[df['Epic status'].str.lower() == 'blocked'].copy()
    blocked_count = len(blocked_df)
    
    # Release analysis
    release_data = []
    if 'Release name' in df.columns:
        for release in df['Release name'].dropna().unique():
            release_epics = df[df['Release name'] == release].copy()
            for idx, row in release_epics.iterrows():
                release_data.append({
                    'Release': release,
                    'Epic ID': row.get('Epic reference #', 'N/A'),
                    'Epic Description': row.get('Epic name', 'N/A'),
                    'Status': row.get('Epic status', 'N/A'),
                    'Epic Link': row.get('Epic URL', 'N/A'),
                    'GitHub URL': row.get('Github enterprise html_url', 'N/A'),
                    'Epic Tags': row.get('Epic tags', 'N/A')
                })
    
    release_df = pd.DataFrame(release_data)
    
    return {
        'total': len(df),
        'status_counts': status_counts,
        'tag_counts': tag_counts,
        'company_count': company_count,
        'company_df': company_df,
        'blocked_count': blocked_count,
        'blocked_df': blocked_df,
        'release_df': release_df,
        'df': df
    }

def create_presentation():
    """Create the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    report_date = datetime.now().strftime('%B %d, %Y')
    add_title_slide(
        prs,
        "2026 Epic Planning & Execution",
        f"IKC & Lineage Analysis Report\n{report_date}"
    )
    
    # Load data
    print("Loading IKC data...")
    ikc_data = analyze_data('aha_list_features_260325061137.xlsx')
    
    print("Loading Lineage data...")
    lineage_data = analyze_data('aha_list_features_260326054547.xlsx')
    
    # Executive Summary
    slide = add_content_slide(prs, "Executive Summary")
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    summary_text = f"""IKC Report:
• Total Epics: {ikc_data['total']}
• Status Categories: {len(ikc_data['status_counts'])}
• Squad Tags: {len(ikc_data['tag_counts'])}
• Company Associations: {ikc_data['company_count']}
• Blocked Epics: {ikc_data['blocked_count']}

Lineage Report:
• Total Epics: {lineage_data['total']}
• Status Categories: {len(lineage_data['status_counts'])}
• Squad Tags: {len(lineage_data['tag_counts'])}
• Company Associations: {lineage_data['company_count']}
• Blocked Epics: {lineage_data['blocked_count']}"""
    
    p = tf.paragraphs[0]
    p.text = summary_text
    p.font.size = Pt(16)
    
    # IKC Section Header
    add_section_header(prs, "IKC 2026 Epic Analysis")
    
    # IKC Status Distribution
    slide = add_content_slide(prs, "IKC: Epic Status Distribution")
    
    chart_data = CategoryChartData()
    chart_data.categories = ikc_data['status_counts'].index.tolist()
    chart_data.add_series('Count', ikc_data['status_counts'].values.tolist())
    
    add_chart_to_slide(
        slide, chart_data, XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(8), Inches(4.5)
    )
    
    # IKC Squad Distribution - ALL squads in table format
    squad_table_data = [['Squad', 'Epic Count']]
    for squad, count in ikc_data['tag_counts'].items():
        squad_table_data.append([squad, str(count)])
    
    # Split into multiple slides if needed (max 20 rows per slide)
    rows_per_slide = 20
    num_slides = math.ceil(len(squad_table_data) / rows_per_slide)
    
    for slide_num in range(num_slides):
        start_idx = slide_num * rows_per_slide
        end_idx = min((slide_num + 1) * rows_per_slide, len(squad_table_data))
        
        if slide_num == 0:
            slide_data = squad_table_data[0:end_idx]
            title = "IKC: All Squad Tags Distribution"
        else:
            slide_data = [squad_table_data[0]] + squad_table_data[start_idx:end_idx]
            title = f"IKC: All Squad Tags Distribution (continued {slide_num + 1})"
        
        add_table_slide(prs, title, slide_data, col_widths=[6, 2])
    
    # IKC Company Association Donut Chart with Table
    if not ikc_data['company_df'].empty:
        slide = add_content_slide(prs, "IKC: Company Association Distribution")
        
        company_counts = ikc_data['company_df']['Company Association'].value_counts().sort_values(ascending=False)
        
        # Create donut chart
        chart_data = CategoryChartData()
        chart_data.categories = company_counts.index.tolist()
        chart_data.add_series('Count', company_counts.values.tolist())
        
        # Add donut chart (left side)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(0.5), Inches(1.8), Inches(4.5), Inches(5),
            chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # Add table (right side)
        table_data = [['Company', 'Count']]
        for company, count in company_counts.items():
            table_data.append([company, str(count)])
        
        rows = len(table_data)
        cols = 2
        
        left = Inches(5.2)
        top = Inches(1.8)
        width = Inches(4.3)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(1.3)
        
        # Fill table
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.size = Pt(11)
                else:
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(9)
        
        # IKC Company Association Table
        company_table_data = [['Epic Reference', 'Epic URL', 'Company Association']]
        for idx, row in ikc_data['company_df'].iterrows():
            epic_ref = str(row.get('Epic reference #', 'N/A'))
            epic_url = str(row.get('Epic URL', 'N/A'))
            company = str(row.get('Company Association', 'N/A'))
            company_table_data.append([epic_ref, epic_url, company])
        
        # Split into multiple slides if needed
        rows_per_slide = 15
        num_slides = math.ceil(len(company_table_data) / rows_per_slide)
        
        for slide_num in range(num_slides):
            start_idx = slide_num * rows_per_slide
            end_idx = min((slide_num + 1) * rows_per_slide, len(company_table_data))
            
            if slide_num == 0:
                slide_data = company_table_data[0:end_idx]
                title = "IKC: Company Association Details"
            else:
                slide_data = [company_table_data[0]] + company_table_data[start_idx:end_idx]
                title = f"IKC: Company Association Details (continued {slide_num + 1})"
            
            add_table_slide(prs, title, slide_data, col_widths=[1.5, 4, 3])
    
    # IKC Blocked Epics
    if not ikc_data['blocked_df'].empty:
        add_blocked_epics_table(prs, "⚠️ IKC: BLOCKED EPICS ⚠️", ikc_data['blocked_df'])
    
    # Lineage Section Header
    add_section_header(prs, "Lineage 2026 Epic Analysis")
    
    # Lineage Status Distribution
    slide = add_content_slide(prs, "Lineage: Epic Status Distribution")
    
    chart_data = CategoryChartData()
    chart_data.categories = lineage_data['status_counts'].index.tolist()
    chart_data.add_series('Count', lineage_data['status_counts'].values.tolist())
    
    add_chart_to_slide(
        slide, chart_data, XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(8), Inches(4.5)
    )
    
    # Lineage Squad Distribution - ALL squads in table format
    squad_table_data = [['Squad', 'Epic Count']]
    for squad, count in lineage_data['tag_counts'].items():
        squad_table_data.append([squad, str(count)])
    
    # Split into multiple slides if needed
    rows_per_slide = 20
    num_slides = math.ceil(len(squad_table_data) / rows_per_slide)
    
    for slide_num in range(num_slides):
        start_idx = slide_num * rows_per_slide
        end_idx = min((slide_num + 1) * rows_per_slide, len(squad_table_data))
        
        if slide_num == 0:
            slide_data = squad_table_data[0:end_idx]
            title = "Lineage: All Squad Tags Distribution"
        else:
            slide_data = [squad_table_data[0]] + squad_table_data[start_idx:end_idx]
            title = f"Lineage: All Squad Tags Distribution (continued {slide_num + 1})"
        
        add_table_slide(prs, title, slide_data, col_widths=[6, 2])
    
    # Lineage Company Association Donut Chart with Table
    if not lineage_data['company_df'].empty:
        slide = add_content_slide(prs, "Lineage: Company Association Distribution")
        
        company_counts = lineage_data['company_df']['Company Association'].value_counts().sort_values(ascending=False)
        
        # Create donut chart
        chart_data = CategoryChartData()
        chart_data.categories = company_counts.index.tolist()
        chart_data.add_series('Count', company_counts.values.tolist())
        
        # Add donut chart (left side)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(0.5), Inches(1.8), Inches(4.5), Inches(5),
            chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # Add table (right side)
        table_data = [['Company', 'Count']]
        for company, count in company_counts.items():
            table_data.append([company, str(count)])
        
        rows = len(table_data)
        cols = 2
        
        left = Inches(5.2)
        top = Inches(1.8)
        width = Inches(4.3)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(1.3)
        
        # Fill table
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.size = Pt(11)
                else:
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(9)
        
        # Lineage Company Association Table
        company_table_data = [['Epic Reference', 'Epic URL', 'Company Association']]
        for idx, row in lineage_data['company_df'].iterrows():
            epic_ref = str(row.get('Epic reference #', 'N/A'))
            epic_url = str(row.get('Epic URL', 'N/A'))
            company = str(row.get('Company Association', 'N/A'))
            company_table_data.append([epic_ref, epic_url, company])
        
        # Split into multiple slides if needed
        rows_per_slide = 15
        num_slides = math.ceil(len(company_table_data) / rows_per_slide)
        
        for slide_num in range(num_slides):
            start_idx = slide_num * rows_per_slide
            end_idx = min((slide_num + 1) * rows_per_slide, len(company_table_data))
            
            if slide_num == 0:
                slide_data = company_table_data[0:end_idx]
                title = "Lineage: Company Association Details"
            else:
                slide_data = [company_table_data[0]] + company_table_data[start_idx:end_idx]
                title = f"Lineage: Company Association Details (continued {slide_num + 1})"
            
            add_table_slide(prs, title, slide_data, col_widths=[1.5, 4, 3])
    
    # Lineage Blocked Epics
    if not lineage_data['blocked_df'].empty:
        add_blocked_epics_table(prs, "⚠️ Lineage: BLOCKED EPICS ⚠️", lineage_data['blocked_df'])
    
    # IKC Release Analysis
    if not ikc_data['release_df'].empty:
        add_section_header(prs, "IKC: Release Analysis")
        
        # Define status colors
        status_colors = {
            'New': RGBColor(135, 206, 235),
            'In Development': RGBColor(255, 215, 0),
            'In Design': RGBColor(221, 160, 221),
            'Dev Complete': RGBColor(144, 238, 144),
            'Shipped': RGBColor(50, 205, 50),
            'Ready for Development': RGBColor(255, 165, 0),
            'Ready for Design': RGBColor(255, 105, 180),
            'Under Consideration': RGBColor(211, 211, 211),
            'Blocked': RGBColor(255, 0, 0),
            'UX Design Delivered': RGBColor(147, 112, 219)
        }
        
        for release in ikc_data['release_df']['Release'].unique():
            release_epics = ikc_data['release_df'][ikc_data['release_df']['Release'] == release]
            
            # Create table data with new columns
            table_data = [['Epic ID', 'Epic Description', 'Status', 'Epic Link', 'GitHub URL', 'Epic Tags']]
            for idx, row in release_epics.iterrows():
                epic_id = str(row.get('Epic ID', 'N/A'))
                epic_desc = str(row.get('Epic Description', 'N/A'))
                status = str(row.get('Status', 'N/A'))
                epic_link = str(row.get('Epic Link', 'N/A'))
                git_url = str(row.get('GitHub URL', 'N/A'))
                tags = str(row.get('Epic Tags', 'N/A'))
                
                # Truncate long fields
                if len(epic_desc) > 30:
                    epic_desc = epic_desc[:27] + '...'
                if len(tags) > 25:
                    tags = tags[:22] + '...'
                
                table_data.append([epic_id, epic_desc, status, epic_link, git_url, tags])
            
            # Split into multiple slides if needed (fewer rows due to more columns)
            rows_per_slide = 15
            num_slides = math.ceil(len(table_data) / rows_per_slide)
            
            for slide_num in range(num_slides):
                start_idx = slide_num * rows_per_slide
                end_idx = min((slide_num + 1) * rows_per_slide, len(table_data))
                
                if slide_num == 0:
                    slide_data = table_data[0:end_idx]
                    title = f"IKC Release: {release}"
                else:
                    slide_data = [table_data[0]] + table_data[start_idx:end_idx]
                    title = f"IKC Release: {release} (continued {slide_num + 1})"
                
                slide = add_content_slide(prs, title)
                
                rows = len(slide_data)
                cols = 6
                
                left = Inches(0.3)
                top = Inches(1.8)
                width = Inches(9.4)
                height = Inches(5)
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column widths
                table.columns[0].width = Inches(0.8)   # Epic ID
                table.columns[1].width = Inches(1.8)   # Epic Description
                table.columns[2].width = Inches(1.2)   # Status
                table.columns[3].width = Inches(2.2)   # Epic Link
                table.columns[4].width = Inches(2.2)   # GitHub URL
                table.columns[5].width = Inches(1.2)   # Epic Tags
                
                # Fill table
                for i, row_data in enumerate(slide_data):
                    for j, cell_data in enumerate(row_data):
                        cell = table.cell(i, j)
                        cell.text = str(cell_data)
                        cell.text_frame.word_wrap = True
                        
                        if i == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                    run.font.size = Pt(9)
                        else:
                            # Color code status column
                            if j == 2:  # Status is now column 2
                                status = slide_data[i][2]
                                color = status_colors.get(status, RGBColor(255, 255, 255))
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = color
                            elif i % 2 == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                            
                            # Make URLs clickable (columns 3 and 4)
                            if j in [3, 4] and str(cell_data) != 'N/A':
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(7)
                                        run.font.color.rgb = RGBColor(0, 0, 255)
                                        run.font.underline = True
                                        run.hyperlink.address = str(cell_data)
                            else:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(7)
    
    # Lineage Release Analysis
    if not lineage_data['release_df'].empty:
        add_section_header(prs, "Lineage: Release Analysis")
        
        # Define status colors
        status_colors = {
            'New': RGBColor(135, 206, 235),
            'In Development': RGBColor(255, 215, 0),
            'In Design': RGBColor(221, 160, 221),
            'Dev Complete': RGBColor(144, 238, 144),
            'Shipped': RGBColor(50, 205, 50),
            'Ready for Development': RGBColor(255, 165, 0),
            'Ready for Design': RGBColor(255, 105, 180),
            'Under Consideration': RGBColor(211, 211, 211),
            'Blocked': RGBColor(255, 0, 0),
            'UX Design Delivered': RGBColor(147, 112, 219)
        }
        
        for release in lineage_data['release_df']['Release'].unique():
            release_epics = lineage_data['release_df'][lineage_data['release_df']['Release'] == release]
            
            # Create table data with new columns
            table_data = [['Epic ID', 'Epic Description', 'Status', 'Epic Link', 'GitHub URL', 'Epic Tags']]
            for idx, row in release_epics.iterrows():
                epic_id = str(row.get('Epic ID', 'N/A'))
                epic_desc = str(row.get('Epic Description', 'N/A'))
                status = str(row.get('Status', 'N/A'))
                epic_link = str(row.get('Epic Link', 'N/A'))
                git_url = str(row.get('GitHub URL', 'N/A'))
                tags = str(row.get('Epic Tags', 'N/A'))
                
                # Truncate long fields
                if len(epic_desc) > 30:
                    epic_desc = epic_desc[:27] + '...'
                if len(tags) > 25:
                    tags = tags[:22] + '...'
                
                table_data.append([epic_id, epic_desc, status, epic_link, git_url, tags])
            
            # Split into multiple slides if needed (fewer rows due to more columns)
            rows_per_slide = 15
            num_slides = math.ceil(len(table_data) / rows_per_slide)
            
            for slide_num in range(num_slides):
                start_idx = slide_num * rows_per_slide
                end_idx = min((slide_num + 1) * rows_per_slide, len(table_data))
                
                if slide_num == 0:
                    slide_data = table_data[0:end_idx]
                    title = f"Lineage Release: {release}"
                else:
                    slide_data = [table_data[0]] + table_data[start_idx:end_idx]
                    title = f"Lineage Release: {release} (continued {slide_num + 1})"
                
                slide = add_content_slide(prs, title)
                
                rows = len(slide_data)
                cols = 6
                
                left = Inches(0.3)
                top = Inches(1.8)
                width = Inches(9.4)
                height = Inches(5)
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set column widths
                table.columns[0].width = Inches(0.8)   # Epic ID
                table.columns[1].width = Inches(1.8)   # Epic Description
                table.columns[2].width = Inches(1.2)   # Status
                table.columns[3].width = Inches(2.2)   # Epic Link
                table.columns[4].width = Inches(2.2)   # GitHub URL
                table.columns[5].width = Inches(1.2)   # Epic Tags
                
                # Fill table
                for i, row_data in enumerate(slide_data):
                    for j, cell_data in enumerate(row_data):
                        cell = table.cell(i, j)
                        cell.text = str(cell_data)
                        cell.text_frame.word_wrap = True
                        
                        if i == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                    run.font.size = Pt(9)
                        else:
                            # Color code status column
                            if j == 2:  # Status is now column 2
                                status = slide_data[i][2]
                                color = status_colors.get(status, RGBColor(255, 255, 255))
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = color
                            elif i % 2 == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                            
                            # Make URLs clickable (columns 3 and 4)
                            if j in [3, 4] and str(cell_data) != 'N/A':
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(7)
                                        run.font.color.rgb = RGBColor(0, 0, 255)
                                        run.font.underline = True
                                        run.hyperlink.address = str(cell_data)
                            else:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(7)
    
    # Comparison Slide
    slide = add_content_slide(prs, "IKC vs Lineage Comparison")
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Total Epics', 'Status Categories', 'Squad Tags', 'Company Assoc.', 'Blocked']
    chart_data.add_series('IKC', [
        ikc_data['total'],
        len(ikc_data['status_counts']),
        len(ikc_data['tag_counts']),
        ikc_data['company_count'],
        ikc_data['blocked_count']
    ])
    chart_data.add_series('Lineage', [
        lineage_data['total'],
        len(lineage_data['status_counts']),
        len(lineage_data['tag_counts']),
        lineage_data['company_count'],
        lineage_data['blocked_count']
    ])
    
    add_chart_to_slide(
        slide, chart_data, XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2), Inches(8), Inches(4.5)
    )
    
    # Save presentation
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f'aha_epic_analysis_combined_{timestamp}.pptx'
    prs.save(filename)
    
    print(f"\n✓ PowerPoint presentation created: {filename}")
    print(f"  - Total slides: {len(prs.slides)}")
    print(f"  - IKC blocked epics: {ikc_data['blocked_count']}")
    print(f"  - Lineage blocked epics: {lineage_data['blocked_count']}")
    return filename

if __name__ == "__main__":
    print("="*70)
    print("Creating PowerPoint Presentation")
    print("="*70)
    
    try:
        filename = create_presentation()
        print("\n" + "="*70)
        print("Presentation created successfully!")
        print("="*70)
        print(f"\nFile: {filename}")
    except Exception as e:
        print(f"\nError: {e}")
        import traceback
        traceback.print_exc()

# Made with Bob
